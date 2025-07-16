from pptx import Presentation
from reportlab.pdfgen import canvas

class PPTPDFConvertisseur:
    def convertir(self, chemin_entree, chemin_sortie):
        # Conversion simple : exporte le texte de chaque slide dans un PDF
        prs = Presentation(chemin_entree)
        c = canvas.Canvas(chemin_sortie)
        y = 800
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            c.drawString(40, y, f"Slide {i+1}:")
            y -= 20
            for line in text:
                c.drawString(60, y, line)
                y -= 15
                if y < 40:
                    c.showPage()
                    y = 800
            y -= 20
            if y < 40:
                c.showPage()
                y = 800
        c.save()